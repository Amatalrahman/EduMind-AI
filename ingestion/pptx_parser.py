"""
PPTX parser using python-pptx.
Treats each slide as one "page", extracts text from all text frames.
"""

from __future__ import annotations

import logging
from pathlib import Path
from pptx import Presentation
from ingestion.pdf_parser import PageContent
from pptx.enum.shapes import MSO_SHAPE_TYPE
from ingestion.text_cleaner import clean_text

logger = logging.getLogger(__name__)


def parse_pptx(file_path: str | Path) -> list[PageContent]:
    """
    Open a PPTX and extract text from each slide.
    Images embedded in slides are extracted as PNG bytes.

    Args:
        file_path: path to the PPTX file on disk.

    Returns:
        List of PageContent objects, one per slide (page_number = slide index, 1-based).
    """
    file_path = Path(file_path)
    prs = Presentation(str(file_path))
    pages: list[PageContent] = []

    logger.info("Parsing PPTX: %s  (%d slides)", file_path.name, len(prs.slides))

    for slide_index, slide in enumerate(prs.slides):
        page_number = slide_index + 1  # 1-indexed

        # ── Text ──────────────────────────────────────────────────────────────
        text_parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(run.text for run in para.runs if run.text.strip())
                    if line.strip():
                        text_parts.append(line.strip())
        text = "\n".join(text_parts)
        text = clean_text(text)

        # ── Images ────────────────────────────────────────────────────────────
        raw_images: list[bytes] = []
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    raw_images.append(image.blob)
                except Exception as exc:
                    logger.warning("Could not extract image on slide %d: %s", page_number, exc)

        pages.append(PageContent(page_number=page_number, text=text, images=raw_images))

    logger.info("Finished parsing PPTX: %d slides extracted", len(pages))
    return pages
